#!/usr/bin/env python3
"""Build the editable Robot Gym investor pitch deck."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "investor"
PPTX_PATH = OUT_DIR / "robot_gym_investor_pitch_deck.pptx"
NOTES_PATH = OUT_DIR / "robot_gym_investor_pitch_deck_notes.md"
ARENA_IMAGE = ROOT / "demo_5" / "web" / "arena-poster.jpg"

W = 13.333
H = 7.5

BG = "080B10"
PANEL = "111722"
PANEL_2 = "171E2A"
WHITE = "F5F7FA"
MUTED = "A6AFBE"
FAINT = "697386"
GRID = "273142"
BLUE = "4A9CFF"
BLUE_2 = "8DC6FF"
RED = "FF4B38"
GREEN = "55D6A5"
YELLOW = "F4C95D"

FONT = "Aptos"
MONO = "Aptos Mono"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def rect(slide, x, y, w, h, fill, radius=False, line=None, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if transparency:
        alpha = OxmlElement("a:alpha")
        alpha.set("val", str(int((100 - transparency) * 1000)))
        shape.fill._xPr.solidFill.srgbClr.append(alpha)
    shape.line.color.rgb = rgb(line or fill)
    if radius:
        shape.adjustments[0] = 0.12
    return shape


def line(slide, x1, y1, x2, y2, color=GRID, width=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x1),
        Inches(y1),
        Inches(max(x2 - x1, 0.01)),
        Inches(max(y2 - y1, 0.01)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def text(
    slide,
    value,
    x,
    y,
    w,
    h,
    *,
    size=18,
    color=WHITE,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
    tracking=None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    if tracking is not None:
        run.font._element.set("spc", str(tracking))
    return box


def rich_text(slide, parts, x, y, w, h, *, size=18, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    for value, color, bold in parts:
        run = p.add_run()
        run.text = value
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_image_cover(slide, path: Path, x, y, w, h, transparency=None):
    with Image.open(path) as im:
        iw, ih = im.size
    image_ratio = iw / ih
    box_ratio = w / h
    if image_ratio > box_ratio:
        visible = box_ratio / image_ratio
        crop = (1 - visible) / 2
        pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        visible = image_ratio / box_ratio
        crop = (1 - visible) / 2
        pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        pic.crop_top = crop
        pic.crop_bottom = crop
    if transparency is not None:
        # python-pptx does not expose picture transparency; use overlays instead.
        rect(slide, x, y, w, h, BG, transparency=transparency)
    return pic


def circle(slide, x, y, d, fill, line_color=None, line_width=1):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = rgb(fill)
    c.line.color.rgb = rgb(line_color or fill)
    c.line.width = Pt(line_width)
    return c


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(BG)
    return slide


def footer(slide, n, section):
    line(slide, 0.48, 7.17, 12.38, 7.18, GRID)
    text(slide, "ROBOT GYM  /  CONFIDENTIAL", 0.5, 7.22, 3.2, 0.18, size=8, color=FAINT, bold=True, font=MONO)
    text(slide, section.upper(), 9.4, 7.22, 2.9, 0.18, size=8, color=FAINT, bold=True, font=MONO, align=PP_ALIGN.RIGHT)
    text(slide, f"{n:02d}", 12.43, 7.22, 0.4, 0.18, size=8, color=WHITE, bold=True, font=MONO, align=PP_ALIGN.RIGHT)


def eyebrow(slide, value, color=BLUE):
    text(slide, value.upper(), 0.55, 0.42, 4.8, 0.25, size=10, color=color, bold=True, font=MONO)


def title(slide, value, y=0.82, size=30, w=11.9):
    text(slide, value, 0.55, y, w, 0.9, size=size, color=WHITE, bold=True)


def pill(slide, value, x, y, w, color=BLUE, text_color=BG):
    rect(slide, x, y, w, 0.32, color, radius=True)
    text(slide, value.upper(), x, y + 0.06, w, 0.14, size=8, color=text_color, bold=True, font=MONO, align=PP_ALIGN.CENTER)


def add_bullet(slide, headline, body, x, y, w, accent=BLUE):
    circle(slide, x, y + 0.08, 0.1, accent)
    text(slide, headline, x + 0.24, y, w - 0.24, 0.28, size=15, bold=True)
    text(slide, body, x + 0.24, y + 0.33, w - 0.24, 0.65, size=11, color=MUTED)


def add_metric(slide, value, label, x, y, w, color=WHITE):
    text(slide, value, x, y, w, 0.58, size=28, color=color, bold=True)
    text(slide, label.upper(), x, y + 0.62, w, 0.38, size=9, color=MUTED, bold=True, font=MONO)


def build_deck():
    if not ARENA_IMAGE.exists():
        raise FileNotFoundError(f"Missing arena image: {ARENA_IMAGE}")
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    props = prs.core_properties
    props.title = "Robot Gym — Investor Pitch"
    props.subject = "Competitive infrastructure for embodied AI"
    props.author = "Kaushik Sivakumar"
    props.keywords = "robotics, embodied AI, Unitree G1, VLGE, MuJoCo, investor"

    # 1 — Cover
    s = blank_slide(prs)
    add_image_cover(s, ARENA_IMAGE, 6.0, 0, 7.333, 7.5)
    rect(s, 5.85, 0, 1.55, 7.5, BG, transparency=25)
    rect(s, 6.0, 0, 7.333, 7.5, BG, transparency=55)
    rect(s, 0.0, 0.0, 0.13, 7.5, BLUE)
    text(s, "ROBOT GYM", 0.62, 0.56, 3.5, 0.35, size=13, color=BLUE, bold=True, font=MONO)
    text(s, "Competitive\ninfrastructure\nfor embodied AI.", 0.62, 1.42, 6.0, 2.55, size=36, bold=True)
    text(
        s,
        "A live 1v1 arena where robot policies compete, humans intervene, and every action becomes evaluation evidence.",
        0.65,
        4.28,
        4.7,
        1.05,
        size=16,
        color=MUTED,
    )
    pill(s, "Investor brief • 2026", 0.65, 6.35, 1.72, BLUE)
    text(s, "KAUSHIK SIVAKUMAR  /  FOUNDER", 2.55, 6.43, 2.75, 0.18, size=8, color=WHITE, bold=True, font=MONO)
    text(s, "ACTUAL PRODUCT ARENA", 10.45, 6.92, 2.22, 0.2, size=8, color=WHITE, bold=True, font=MONO, align=PP_ALIGN.RIGHT)

    # 2 — Problem
    s = blank_slide(prs)
    eyebrow(s, "The problem", RED)
    title(s, "Robotics demos prove possibility.\nThey rarely prove reliability.", size=31)
    text(s, "The missing layer is not another model. It is a credible way to test models under pressure.", 0.58, 2.5, 7.7, 0.55, size=16, color=MUTED)
    line(s, 0.58, 3.35, 12.18, 3.36, GRID)
    items = [
        ("01", "Success hides failure", "Polished demos compress retries, recoveries, latency and perception misses into one highlight."),
        ("02", "Benchmarks lack stakes", "Static leaderboards reveal scores—not how policies behave in a live, adversarial scenario."),
        ("03", "Hardware iteration is costly", "Every physical test consumes operator time, robot time and safety margin."),
    ]
    for i, (n, h, b) in enumerate(items):
        x = 0.58 + i * 4.13
        text(s, n, x, 3.72, 0.45, 0.28, size=10, color=RED, bold=True, font=MONO)
        text(s, h, x, 4.22, 3.65, 0.48, size=19, bold=True)
        text(s, b, x, 4.9, 3.55, 1.03, size=12, color=MUTED)
    rect(s, 0.58, 6.32, 12.18, 0.52, PANEL, radius=True)
    rich_text(s, [("Opportunity: ", BLUE, True), ("turn evaluation into something rigorous enough for labs—and compelling enough for an audience.", WHITE, False)], 0.84, 6.46, 11.65, 0.25, size=13)
    footer(s, 2, "Problem")

    # 3 — Product
    s = blank_slide(prs)
    add_image_cover(s, ARENA_IMAGE, 0, 0, 7.15, 7.5)
    rect(s, 0, 0, 7.15, 7.5, BG, transparency=68)
    rect(s, 6.65, 0, 0.75, 7.5, BG, transparency=20)
    text(s, "THE PRODUCT", 0.55, 0.43, 2.5, 0.25, size=10, color=BLUE, bold=True, font=MONO)
    text(s, "A live arena for\nrobot intelligence.", 0.55, 0.98, 5.4, 1.35, size=30, bold=True)
    text(s, "Same task. Same physics. Different minds.", 0.57, 2.61, 4.7, 0.35, size=16, color=WHITE, bold=True)
    text(s, "Policies and people race Unitree G1 robots to acquire, carry and deliver a payload.", 0.57, 3.05, 4.95, 0.72, size=13, color=MUTED)
    modes = [("AI", "vs", "AI"), ("AI", "vs", "HUMAN"), ("HUMAN", "vs", "HUMAN")]
    for i, (a, v, b) in enumerate(modes):
        y = 4.18 + i * 0.68
        rect(s, 0.58, y, 4.55, 0.46, PANEL, radius=True)
        text(s, a, 0.82, y + 0.12, 1.25, 0.16, size=9, color=BLUE_2, bold=True, font=MONO)
        text(s, v.upper(), 2.17, y + 0.12, 0.55, 0.16, size=9, color=FAINT, bold=True, font=MONO, align=PP_ALIGN.CENTER)
        text(s, b, 3.0, y + 0.12, 1.75, 0.16, size=9, color=RED if "HUMAN" in b else BLUE_2, bold=True, font=MONO, align=PP_ALIGN.RIGHT)
    text(s, "LIVE PRODUCT CAPTURE  /  VLGE × MUJOCO", 0.58, 6.72, 4.5, 0.2, size=8, color=WHITE, bold=True, font=MONO)
    text(s, "What the customer gets", 7.55, 0.62, 4.8, 0.35, size=12, color=BLUE, bold=True, font=MONO)
    features = [
        ("Watch", "A shared, browser-viewable match with live state and rationales."),
        ("Control", "Gamepad or keyboard input beside model-driven agents."),
        ("Measure", "Trajectories, events, commands, falls and recovery attempts."),
        ("Replay", "An evidence package for debugging, comparison and trust."),
    ]
    for i, (h, b) in enumerate(features):
        add_bullet(s, h, b, 7.55, 1.35 + i * 1.3, 4.75, BLUE if i < 2 else GREEN)
    pill(s, "Built and running", 7.55, 6.57, 1.45, GREEN, BG)
    footer(s, 3, "Product")

    # 4 — Architecture
    s = blank_slide(prs)
    eyebrow(s, "How it works")
    title(s, "One arena. A clean boundary between\nreasoning, control and evidence.", size=29)
    layers = [
        ("01", "VLGE WORLD", "Scenario, visuals, audience", BLUE),
        ("02", "PERCEPTION", "Delayed camera-grounded state", BLUE_2),
        ("03", "PLAYER", "Gemini • custom model • human", YELLOW),
        ("04", "GUARDED SKILLS", "Navigate • grasp • release • reset", RED),
        ("05", "G1 CONTROL", "Official locomotion policy + SDK-shaped channel", GREEN),
        ("06", "EVIDENCE", "Trajectory • events • outcome • replay", WHITE),
    ]
    y = 2.55
    for i, (n, h, b, c) in enumerate(layers):
        x = 0.56 + i * 2.08
        rect(s, x, y, 1.78, 2.15, PANEL if i % 2 == 0 else PANEL_2, radius=True, line=GRID)
        text(s, n, x + 0.18, y + 0.2, 0.35, 0.2, size=9, color=c, bold=True, font=MONO)
        circle(s, x + 1.35, y + 0.18, 0.16, c)
        text(s, h, x + 0.18, y + 0.75, 1.4, 0.42, size=12, color=WHITE, bold=True, font=MONO)
        text(s, b, x + 0.18, y + 1.27, 1.42, 0.64, size=9, color=MUTED)
        if i < len(layers) - 1:
            text(s, "→", x + 1.81, y + 0.92, 0.28, 0.3, size=18, color=FAINT, bold=True, align=PP_ALIGN.CENTER)
    rect(s, 0.56, 5.22, 12.18, 1.0, PANEL, radius=True)
    pill(s, "Key design choice", 0.82, 5.52, 1.42, BLUE)
    text(s, "Models choose grounded intent; the safety and control layer owns executable motion.", 2.52, 5.42, 9.6, 0.44, size=16, color=WHITE, bold=True)
    text(s, "This makes model comparison possible without pretending a language model is a motor controller.", 2.53, 5.9, 9.25, 0.28, size=10, color=MUTED)
    footer(s, 4, "System")

    # 5 — Why now
    s = blank_slide(prs)
    eyebrow(s, "Why now", GREEN)
    title(s, "Embodied AI is moving from demos\nto an evaluation economy.", size=30)
    metrics = [
        ("$38B", "Projected humanoid robot market by 2035", "Goldman Sachs Research", BLUE),
        ("542K", "Industrial robots installed in 2024", "International Federation of Robotics", GREEN),
        ("+42%", "Growth in RaaS fleets in 2024", "International Federation of Robotics", YELLOW),
    ]
    for i, (v, label, source, color) in enumerate(metrics):
        x = 0.58 + i * 4.08
        line(s, x, 2.42, x + 3.5, 2.43, color)
        add_metric(s, v, label, x, 2.72, 3.55, color)
        text(s, source, x, 3.92, 3.48, 0.35, size=9, color=FAINT)
    rect(s, 0.58, 4.68, 12.18, 1.25, PANEL, radius=True)
    text(s, "MODEL SHIFT", 0.84, 5.0, 1.15, 0.18, size=9, color=RED, bold=True, font=MONO)
    text(s, "Embodied reasoning is now API-accessible—planning and spatial understanding can be evaluated across models.", 2.22, 4.92, 9.82, 0.54, size=15, bold=True)
    text(s, "Robot Gym is the neutral arena between rapidly changing models and slow, expensive hardware.", 2.23, 5.49, 9.6, 0.3, size=10, color=MUTED)
    text(s, "Sources: Goldman Sachs Research (2024); IFR World Robotics 2025; Google DeepMind Gemini Robotics-ER.", 0.58, 6.58, 10.9, 0.23, size=8, color=FAINT)
    footer(s, 5, "Market")

    # 6 — Wedge and flywheel
    s = blank_slide(prs)
    eyebrow(s, "The wedge", YELLOW)
    title(s, "Competition turns evaluation\ninto a product.", size=31)
    text(s, "A race is legible to an audience—and unforgiving to a policy.", 0.58, 2.35, 6.5, 0.42, size=16, color=MUTED)
    stages = [
        ("01", "RUN", "Models and humans enter a shared task."),
        ("02", "MEASURE", "Every action becomes structured evidence."),
        ("03", "COMPARE", "Outcomes reveal capability and failure modes."),
        ("04", "IMPROVE", "Teams iterate policies, prompts and controllers."),
    ]
    for i, (n, h, b) in enumerate(stages):
        x = 0.58 + i * 3.08
        circle(s, x, 3.4, 0.66, BLUE if i < 2 else RED, line_color=WHITE)
        text(s, n, x, 3.61, 0.66, 0.16, size=8, color=BG, bold=True, font=MONO, align=PP_ALIGN.CENTER)
        if i < 3:
            line(s, x + 0.75, 3.71, x + 2.88, 3.73, GRID)
        text(s, h, x, 4.34, 2.4, 0.33, size=14, bold=True, font=MONO)
        text(s, b, x, 4.8, 2.46, 0.8, size=10, color=MUTED)
    rect(s, 0.58, 6.05, 12.18, 0.67, PANEL_2, radius=True)
    rich_text(s, [("Flywheel: ", YELLOW, True), ("more scenarios → more matches → better failure data → stronger benchmarks → more teams.", WHITE, False)], 0.88, 6.25, 11.42, 0.28, size=13)
    footer(s, 6, "Wedge")

    # 7 — Built
    s = blank_slide(prs)
    eyebrow(s, "Execution proof", GREEN)
    title(s, "The core arena is already built.", size=31)
    pill(s, "Shipped", 10.94, 0.48, 1.18, GREEN, BG)
    columns = [
        (
            "MATCH",
            [
                "AI vs AI",
                "AI vs human",
                "Human vs human",
                "Rematch + idle player",
            ],
        ),
        (
            "CONTROL",
            [
                "Browser gamepad + keyboard",
                "Official G1 locomotion policy",
                "500 Hz physics / 50 Hz channel",
                "Guarded grasp, release + reset",
            ],
        ),
        (
            "EVALUATION",
            [
                "Camera-grounded perception",
                "Latency, dropout + noise",
                "Custom HTTP policy adapter",
                "Evidence and replay package",
            ],
        ),
    ]
    for i, (heading, bullets) in enumerate(columns):
        x = 0.58 + i * 4.08
        rect(s, x, 2.15, 3.62, 3.65, PANEL if i != 1 else PANEL_2, radius=True, line=GRID)
        text(s, heading, x + 0.25, 2.43, 2.65, 0.25, size=10, color=[BLUE, RED, GREEN][i], bold=True, font=MONO)
        for j, bullet in enumerate(bullets):
            circle(s, x + 0.27, 3.06 + j * 0.62, 0.08, [BLUE, RED, GREEN][i])
            text(s, bullet, x + 0.51, 3.0 + j * 0.62, 2.78, 0.35, size=11, color=WHITE if j == 0 else MUTED)
    text(s, "DEPLOYMENT", 0.58, 6.25, 1.2, 0.2, size=9, color=FAINT, bold=True, font=MONO)
    text(s, "Mac development  →  RunPod GPU host  →  browser spectators and players", 1.92, 6.18, 8.2, 0.34, size=14, bold=True)
    text(s, "Current product scope; not a claim of real-hardware certification.", 9.3, 6.27, 3.45, 0.18, size=8, color=FAINT, align=PP_ALIGN.RIGHT)
    footer(s, 7, "Proof")

    # 8 — Business model
    s = blank_slide(prs)
    eyebrow(s, "Business model", BLUE)
    title(s, "Land as evaluation infrastructure.\nExpand through data and events.", size=30)
    pill(s, "Proposed", 10.78, 0.48, 1.34, YELLOW, BG)
    offerings = [
        ("PLATFORM", "Annual team / enterprise subscription", "Private arenas, model adapters, benchmark management and collaboration."),
        ("COMPUTE", "Metered simulation and replay", "Hosted matches, accelerated scenario runs and evidence retention."),
        ("PROGRAMS", "Benchmarks, tournaments and services", "Custom evaluation suites, partner challenges and failure analysis."),
    ]
    for i, (tag, headline, body) in enumerate(offerings):
        y = 2.33 + i * 1.28
        text(s, f"0{i+1}", 0.6, y + 0.08, 0.4, 0.2, size=9, color=BLUE, bold=True, font=MONO)
        text(s, tag, 1.14, y + 0.06, 1.18, 0.2, size=9, color=FAINT, bold=True, font=MONO)
        text(s, headline, 2.58, y, 4.02, 0.35, size=16, bold=True)
        text(s, body, 7.0, y, 5.1, 0.57, size=11, color=MUTED)
        line(s, 0.59, y + 0.83, 12.2, y + 0.84, GRID)
    rect(s, 0.58, 6.28, 12.18, 0.48, PANEL, radius=True)
    text(s, "Initial buyer: robotics labs and embodied-AI teams that need repeatable cross-model evaluation.", 0.84, 6.42, 11.45, 0.2, size=12, color=WHITE, bold=True)
    footer(s, 8, "Business")

    # 9 — Go to market
    s = blank_slide(prs)
    eyebrow(s, "Go to market", RED)
    title(s, "Start with teams that feel the pain.\nEarn the audience after the signal.", size=30)
    phases = [
        ("NOW", "DESIGN PARTNERS", "Robotics labs + model teams", "Install the arena, define one decision-critical benchmark, measure iteration time."),
        ("NEXT", "PUBLIC CHALLENGE", "Developers + research community", "Publish a task season, verified match evidence and a model-neutral leaderboard."),
        ("SCALE", "PRIVATE ARENAS", "Enterprises + robot OEMs", "Sell proprietary scenarios, access control, fleet evaluation and hosted compute."),
    ]
    for i, (when, phase, buyer, action) in enumerate(phases):
        x = 0.58 + i * 4.08
        rect(s, x, 2.42, 3.62, 3.62, PANEL if i < 2 else PANEL_2, radius=True, line=GRID)
        pill(s, when, x + 0.24, 2.68, 0.77, [BLUE, YELLOW, RED][i], BG)
        text(s, phase, x + 0.24, 3.38, 3.08, 0.33, size=13, bold=True, font=MONO)
        text(s, buyer, x + 0.24, 4.08, 3.03, 0.48, size=16, bold=True)
        text(s, action, x + 0.24, 4.82, 3.02, 0.84, size=10, color=MUTED)
    text(s, "North-star metric", 0.6, 6.42, 1.6, 0.2, size=9, color=FAINT, bold=True, font=MONO)
    text(s, "Verified evaluation runs per active team", 2.25, 6.32, 4.8, 0.35, size=15, color=WHITE, bold=True)
    footer(s, 9, "GTM")

    # 10 — Defensibility
    s = blank_slide(prs)
    eyebrow(s, "Defensibility", YELLOW)
    title(s, "The moat compounds across four layers.", size=31)
    moat = [
        ("01", "SCENARIO GRAPH", "A growing library of tasks, resets, constraints and win conditions."),
        ("02", "FAILURE DATA", "Comparable traces of perception, reasoning, control and recovery."),
        ("03", "CONTROL BOUNDARY", "Hardware-shaped channels that make policy claims more honest."),
        ("04", "COMPETITION NETWORK", "Models, teams, spectators and partners improve the benchmark together."),
    ]
    for i, (n, h, b) in enumerate(moat):
        y = 2.16 + i * 1.05
        rect(s, 0.58, y, 0.62, 0.68, [BLUE, GREEN, RED, YELLOW][i], radius=True)
        text(s, n, 0.58, y + 0.23, 0.62, 0.16, size=8, color=BG, bold=True, font=MONO, align=PP_ALIGN.CENTER)
        text(s, h, 1.48, y + 0.05, 2.62, 0.3, size=13, bold=True, font=MONO)
        text(s, b, 4.25, y + 0.03, 7.7, 0.44, size=12, color=MUTED)
        line(s, 1.48, y + 0.74, 12.2, y + 0.75, GRID)
    rect(s, 0.58, 6.45, 12.18, 0.33, BLUE, radius=True)
    text(s, "Each match improves the product even when the underlying models change.", 0.58, 6.53, 12.18, 0.17, size=10, color=BG, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 10, "Moat")

    # 11 — Roadmap
    s = blank_slide(prs)
    eyebrow(s, "Roadmap", GREEN)
    title(s, "From one compelling race\nto the evaluation layer for embodied AI.", size=30)
    milestones = [
        ("BUILT", "G1 DELIVERY ARENA", ["Three match modes", "Model + human control", "RunPod deployment", "Evidence capture"], GREEN),
        ("12 MONTHS", "HOSTED BENCHMARK SUITE", ["Accounts + private arenas", "Multi-task scenario library", "Verified leaderboards", "Design-partner pilots"], BLUE),
        ("24 MONTHS", "HARDWARE-BACKED NETWORK", ["Real-unit adapter validation", "Partner robot profiles", "Tournament seasons", "Enterprise evaluation API"], RED),
    ]
    for i, (when, head, bullets, color) in enumerate(milestones):
        x = 0.58 + i * 4.08
        line(s, x, 2.55, x + 3.6, 2.57, color)
        text(s, when, x, 2.83, 2.3, 0.22, size=9, color=color, bold=True, font=MONO)
        text(s, head, x, 3.31, 3.42, 0.58, size=16, bold=True)
        for j, bullet in enumerate(bullets):
            circle(s, x, 4.34 + j * 0.46, 0.07, color)
            text(s, bullet, x + 0.22, 4.27 + j * 0.46, 3.16, 0.3, size=10, color=MUTED if j else WHITE)
    text(s, "Milestones are proposed and subject to customer discovery, funding and hardware partner access.", 0.58, 6.56, 10.65, 0.21, size=8, color=FAINT)
    footer(s, 11, "Roadmap")

    # 12 — Ask
    s = blank_slide(prs)
    rect(s, 0, 0, 0.13, 7.5, RED)
    text(s, "THE ASK", 0.62, 0.55, 1.6, 0.25, size=10, color=RED, bold=True, font=MONO)
    text(s, "Help turn robot evaluation\ninto a live, trusted market.", 0.62, 1.25, 8.7, 1.4, size=32, bold=True)
    text(s, "Seeking pre-seed capital and design partners.", 0.64, 2.98, 6.8, 0.42, size=18, color=BLUE_2, bold=True)
    use = [
        ("PRODUCT", "Hosted platform, scenario tooling and evidence UX"),
        ("ROBOTICS", "Benchmark design, controls validation and hardware bridge"),
        ("MARKET", "Design partners, flagship challenge and community"),
    ]
    for i, (h, b) in enumerate(use):
        y = 4.02 + i * 0.72
        text(s, f"0{i+1}", 0.65, y + 0.04, 0.35, 0.2, size=9, color=FAINT, bold=True, font=MONO)
        text(s, h, 1.24, y, 1.35, 0.25, size=10, color=WHITE, bold=True, font=MONO)
        text(s, b, 2.8, y, 4.6, 0.38, size=11, color=MUTED)
    add_image_cover(s, ARENA_IMAGE, 8.22, 0.78, 4.55, 5.38)
    rect(s, 8.22, 0.78, 4.55, 5.38, BG, transparency=78)
    rect(s, 8.52, 4.86, 3.95, 0.95, BG, radius=True, line=GRID, transparency=8)
    text(s, "KAUSHIK SIVAKUMAR", 8.78, 5.09, 3.4, 0.23, size=12, bold=True)
    text(s, "FOUNDER  /  ROBOT GYM", 8.78, 5.45, 3.4, 0.18, size=8, color=BLUE, bold=True, font=MONO)
    text(s, "github.com/KaushikSiva/robot-gym", 0.65, 6.65, 4.2, 0.2, size=10, color=WHITE, bold=True)
    text(s, "LET’S BUILD THE ARENA.", 8.22, 6.62, 4.55, 0.25, size=11, color=RED, bold=True, font=MONO, align=PP_ALIGN.RIGHT)
    footer(s, 12, "Contact")

    prs.save(PPTX_PATH)
    NOTES_PATH.write_text(
        """# Robot Gym investor pitch — presenter notes

## Positioning

Robot Gym is competitive infrastructure for embodied AI: a live arena where robot
policies and human operators solve the same physical task under shared constraints.
The match is the interface; structured evaluation evidence is the product.

## Suggested 8–10 minute talk track

1. **Cover** — Models are improving quickly, but robotics teams still lack a trusted,
   legible way to compare embodied behavior.
2. **Problem** — A polished demo can hide retries, latency and recovery. Hardware-only
   iteration is slow and expensive.
3. **Product** — Robot Gym makes evaluation watchable: AI vs AI, AI vs human, or human
   vs human in a shared Unitree G1 arena.
4. **System** — Models decide grounded intent. Guarded skills, locomotion and an
   SDK-shaped channel own execution. Evidence is captured end-to-end.
5. **Why now** — Robot deployment is growing, humanoid investment is accelerating,
   and embodied reasoning is now available through model APIs.
6. **Wedge** — Competition creates stakes, repeated runs and clear outcomes. Those
   matches become a high-value failure dataset and benchmark suite.
7. **Proof** — The core Demo 5 product is built and deployable on Mac and RunPod.
   Be precise: this is simulation evidence, not real-hardware certification.
8. **Business** — Proposed revenue comes from platform subscriptions, metered compute,
   and custom benchmarks/tournaments.
9. **GTM** — Begin with design partners who already spend heavily on evaluation;
   add a public challenge only after the benchmark signal is credible.
10. **Moat** — Scenarios, failure traces, control boundaries and the participant
    network compound even as foundation models change.
11. **Roadmap** — Expand from one G1 race to a hosted suite, then validate real-unit
    adapters with hardware partners.
12. **Ask** — Pre-seed capital plus 3–5 design partners to turn the working arena into
    a hosted evaluation product.

## Claims discipline

- Slides marked **SHIPPED** describe repository capabilities.
- Slides marked **PROPOSED** describe the business plan, not current revenue.
- Do not describe the current simulation as proof that the system is certified or
  safe for unsupervised operation on a real Unitree G1.

## Sources

- Goldman Sachs Research, “The global market for robots could reach $38 billion by
  2035” (2024): https://www.goldmansachs.com/insights/articles/the-global-market-for-robots-could-reach-38-billion-by-2035.html
- International Federation of Robotics, World Robotics 2025:
  https://ifr.org/worldrobotics/report-2025
- Google DeepMind, “Gemini Robotics brings AI into the physical world” (2025):
  https://deepmind.google/blog/gemini-robotics-brings-ai-into-the-physical-world/
- Google DeepMind, “Gemini Robotics-ER 1.6” (2026):
  https://deepmind.google/blog/gemini-robotics-er-1-6/
- Unitree G1 product page: https://www.unitree.com/g1/
- VLGE AI: https://vlge.com/ai
""",
        encoding="utf-8",
    )
    print(PPTX_PATH)
    print(NOTES_PATH)


if __name__ == "__main__":
    build_deck()
